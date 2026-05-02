#!/usr/bin/env python3
"""Generate the Resonance Installation Deck — for curators and venue programmers.

Mirrors the visual system of the YC vision deck (build-deck.py) — same
dimensions, dark theme, purple accent, Calibri / Consolas substitutions
for Matter / Geist — but with content tuned for curators of immersive
listening rooms (Envelop SF, planetariums, sound-art venues, festivals).

Source narrative lives in docs/installation-brief.md. This script is the
slide-format layer.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Visual system (matches build-deck.py) ─────────────────────────────
W = Inches(13.333)  # 1440px @ 108dpi → 16:9
H = Inches(7.5)
BG = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

# White at decreasing opacities, approximated as solid grey on black
W90 = WHITE
W75 = RGBColor(0xBF, 0xBF, 0xBF)
W60 = RGBColor(0x99, 0x99, 0x99)
W50 = RGBColor(0x80, 0x80, 0x80)
W45 = RGBColor(0x73, 0x73, 0x73)
W40 = RGBColor(0x66, 0x66, 0x66)
W30 = RGBColor(0x4D, 0x4D, 0x4D)
W25 = RGBColor(0x40, 0x40, 0x40)
W20 = RGBColor(0x33, 0x33, 0x33)
W15 = RGBColor(0x26, 0x26, 0x26)

P60 = RGBColor(0xA7, 0x78, 0xF2)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=14,
             color=WHITE, bold=False, italic=False, font_name="Calibri",
             alignment=PP_ALIGN.LEFT, line_spacing=None, spacing_after=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if spacing_after is not None:
        p.space_after = Pt(spacing_after)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_name="Calibri"):
    """lines = list of (text, font_size, color, bold, italic, line_spacing?, spacing_after?)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, fs, clr, bld, ital = line[0], line[1], line[2], line[3], line[4]
        ls = line[5] if len(line) > 5 else None
        sa = line[6] if len(line) > 6 else None
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.italic = ital
        p.font.name = font_name
        if ls:
            p.line_spacing = Pt(ls)
        if sa is not None:
            p.space_after = Pt(sa)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


# ── Margins ────────────────────────────────────────────────────────────
LM = Inches(0.9)
TM = Inches(0.67)
CW = Inches(11.5)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 01 — Title
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(2.5), Inches(1), Inches(0.4), "⌇", 24, W60)

add_text(s, LM, Inches(3.0), Inches(11), Inches(1.0),
         "Resonance", 52, WHITE, bold=True)

add_text(s, LM, Inches(4.1), Inches(11), Inches(1.0),
         "A contemplative listening room.\nGenerative audiovisual, music-led, slow time.",
         18, W50, line_spacing=28)

add_text(s, LM, Inches(6.5), Inches(8), Inches(0.3),
         "INSTALLATION BRIEF    |    KAREL BARNOSKI    |    MAY 2026",
         9, W25, font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 02 — What it is
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.4), Inches(3), Inches(0.25),
         "WHAT IT IS", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.75), Inches(11), Inches(1.6),
         "A generative audiovisual instrument.\nMusic-led. Slow time.",
         44, WHITE, bold=True, line_spacing=52)

add_multiline(s, LM, Inches(4.0), Inches(11), Inches(2.5), [
    ("A composed piece of music drives real-time audio-reactive shaders and "
     "AI-curated imagery, producing a 20–40 minute slow-moving visual landscape "
     "that never repeats verbatim.", 15, W60, False, False, 24, 14),
    ("Originally built for personal listening. Scales naturally to a shared "
     "room: same engine, same patience, just larger and quieter.",
     15, W60, False, False, 24),
])


# ══════════════════════════════════════════════════════════════════════
# SLIDE 03 — The Room
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.0), Inches(3), Inches(0.25),
         "THE ROOM", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.35), Inches(11), Inches(1.6),
         "Planetarium vibes.\nA dark space, audience reclining, eyes up.",
         42, WHITE, bold=True, line_spacing=50)

specs = [
    ("Space", "12×12 ft to 30×40 ft. Dome geometry welcome."),
    ("Floor", "Mats and pillows. No chairs."),
    ("Display", "Single overhead projection, dome, or large wall display."),
    ("Audio", "4-corner speaker array. Stereo as a baseline."),
    ("Flow", "Visitors enter and leave on their own time. No timed entry. No headphones. No narration."),
]
for i, (label, desc) in enumerate(specs):
    y = Inches(4.3 + i * 0.45)
    add_text(s, LM, y, Inches(1.6), Inches(0.4),
             label.upper(), 9, PURPLE, font_name="Consolas")
    add_text(s, LM + Inches(1.7), y, Inches(9.5), Inches(0.4),
             desc, 13, W60, line_spacing=20)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 04 — The Piece
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.0), Inches(3), Inches(0.25),
         "THE PIECE", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.35), Inches(11), Inches(1.2),
         "A single ~30 minute Path.", 42, WHITE, bold=True)

add_text(s, LM, Inches(2.7), Inches(11), Inches(0.9),
         "Three to five composed tracks, threaded into one continuous arc.",
         16, W45, line_spacing=24)

# Phase chart
phases = [("01", "Threshold", "Open in near-silence. Single tone. Near-black field."),
          ("02", "Expansion", "Imagery surfaces. Music builds. Caustics, mandalas, figures."),
          ("03", "Apex", "Full motion and color. The room is loud and alive."),
          ("04", "Return", "Settle back toward stillness. Cycle resolves and begins again.")]
card_w = Inches(2.7)
gap = Inches(0.18)
for i, (num, label, desc) in enumerate(phases):
    x = LM + (card_w + gap) * i
    y = Inches(4.4)
    add_rect(s, x, y, card_w, Inches(2.3), border_color=W15)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(0.25),
             num, 10, PURPLE, font_name="Consolas")
    add_text(s, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.4), Inches(0.4),
             label, 17, WHITE, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(1.1), card_w - Inches(0.4), Inches(1.2),
             desc, 11, W45, line_spacing=17)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 04.5 — Shape of the work (image: Welcome Home path landing)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(0.55), Inches(6), Inches(0.25),
         "SHAPE OF THE WORK", 9, PURPLE, font_name="Consolas")

# Full-bleed image, slight inset so the dark page bleeds around it.
still_path = os.path.join(os.path.dirname(__file__), "..", "public",
                          "installation-stills", "01-welcome-home-path.png")
if os.path.exists(still_path):
    s.shapes.add_picture(still_path, Inches(1.3), Inches(1.0),
                         width=Inches(10.7), height=Inches(5.6))

add_text(s, Inches(0), Inches(6.85), W, Inches(0.3),
         "Welcome Home — thirteen tracks, one ~30-minute path through the album",
         11, W30, italic=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 05 — What a visitor experiences (the moment)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(0.9), Inches(6), Inches(0.25),
         "WHAT A VISITOR EXPERIENCES", 9, PURPLE, font_name="Consolas")

# A wide italic narrative paragraph centered vertically.
add_text(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(4.0),
         "Enters a dark room. Lies back. The first phase opens with a single low tone "
         "and a near-black field — slow caustics, the faint hint of a mandalic form. "
         "Imagery surfaces gradually as the music builds: water, light, figures that "
         "resolve and dissolve. Around the apex the room is full of motion and color; "
         "then everything settles back toward stillness. The visitor leaves when ready.",
         24, W75, italic=True, line_spacing=38)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 06 — References
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.0), Inches(3), Inches(0.25),
         "REFERENCES", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.35), Inches(11), Inches(1.0),
         "A small canon.", 42, WHITE, bold=True)

refs = [
    ("Brian Eno",
     "77 Million Paintings",
     "Generative visual companion to ambient music. Proof that slow time finds an audience."),
    ("James Turrell",
     "Skyspaces (1974— )",
     "The dark contemplative room as primary form. No narrative scaffolding. Pure perception."),
    ("Ryoji Ikeda",
     "test pattern  ·  data.matrix",
     "Sound and light as total sensorium. Room-scale, durational, austere."),
]
card_w = Inches(3.65)
gap = Inches(0.27)
for i, (artist, work, desc) in enumerate(refs):
    x = LM + (card_w + gap) * i
    y = Inches(3.3)
    add_rect(s, x, y, card_w, Inches(3.2), border_color=W15)
    add_text(s, x + Inches(0.25), y + Inches(0.3), card_w - Inches(0.5), Inches(0.4),
             artist, 18, WHITE, bold=True)
    add_text(s, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.45),
             work, 12, P60, italic=True, font_name="Consolas")
    add_text(s, x + Inches(0.25), y + Inches(1.55), card_w - Inches(0.5), Inches(1.5),
             desc, 12, W50, line_spacing=18)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 07 — Where it differs
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.4), Inches(3), Inches(0.25),
         "WHERE IT DIFFERS", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.75), Inches(11), Inches(2.4),
         "Composed first.\nAI second.",
         60, WHITE, bold=True, line_spacing=72)

add_multiline(s, LM, Inches(5.0), Inches(11), Inches(2.5), [
    ("The AI-image installation has become a genre — Anadol-flavored generative "
     "spectacle, large-format, often impressive, often hollow.",
     15, W50, False, False, 24, 14),
    ("Resonance is the inverse. Each track is a written piece of music with intent. "
     "The visuals serve the music, not the other way around. AI sits inside the "
     "journey as a tool. It is not the headline.",
     15, W60, False, False, 24),
])


# ══════════════════════════════════════════════════════════════════════
# SLIDE 08 — Why now
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.0), Inches(3), Inches(0.25),
         "WHY NOW", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.35), Inches(11), Inches(1.6),
         "Slow attention is having a moment.",
         42, WHITE, bold=True, line_spacing=52)

evidence = [
    ("Hiroshi Yoshimura's Music for Nine Post Cards hit the Billboard charts in 2024 — "
     "fifty years after release. Ambient music is mainstream again."),
    ("Planetariums are reclaiming a cultural foothold. Hayden, Adler, and Morrison have "
     "all renovated and reopened their domes for contemporary programming."),
    ("Sleep No More-era audiences trained themselves on long-form, low-direction "
     "experiences. The patience is back."),
]
for i, line in enumerate(evidence):
    y = Inches(3.7 + i * 1.05)
    # Purple bullet
    add_rect(s, LM, y + Inches(0.18), Inches(0.06), Inches(0.06), fill_color=PURPLE)
    add_text(s, LM + Inches(0.3), y, Inches(11), Inches(0.9),
             line, 14, W60, line_spacing=22)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 09 — Who
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "WHO", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(11), Inches(0.7),
         "Karel Barnoski", 42, WHITE, bold=True)

add_text(s, LM, Inches(1.7), Inches(11), Inches(0.4),
         "Pianist  ·  Composer  ·  Designer",
         14, P60, italic=True)

bio_left = [
    ("Pianist and composer. Released Welcome Home, a solo piano album self-produced "
     "during quarantine — the source material for many of the journeys in Resonance.",
     12, W60, False, False, 19, 14),
    ("Design Director at Workday. 15+ years leading UX for enterprise products at "
     "Workday, GE, and Kodak. BFA Illustration and MFA Computer Graphics from RIT, "
     "with classical training from Barnstone Studios.",
     12, W50, False, False, 19),
]
add_multiline(s, LM, Inches(2.5), Inches(5.3), Inches(3.0), bio_left)

bio_right = [
    ("Founder of 2octave, an audio design studio specializing in product sonification. "
     "Award-winning sounds for products from digital cameras to kitchen appliances. "
     "Published in UX Magazine on the science of sound in product design.",
     12, W50, False, False, 19, 14),
    ("Resonance is the current shape of a multi-year practice translating composed "
     "music into shared visual experience. The installation is the natural next room.",
     12, W60, False, False, 19),
]
add_multiline(s, LM + Inches(6.0), Inches(2.5), Inches(5.5), Inches(3.0), bio_right)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — The ask + contact
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, Inches(0), Inches(1.5), W, Inches(0.3),
         "THE ASK", 9, PURPLE, font_name="Consolas",
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(2.0), W, Inches(1.5),
         "Let's talk.", 64, WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(3.4), W, Inches(2.5),
         "A one-night event. A multi-week run. A residency.\n"
         "A slot in an existing program. A conversation.\n"
         "Whatever shape the fit takes.",
         18, W50, alignment=PP_ALIGN.CENTER, line_spacing=30)

add_text(s, Inches(0), Inches(5.6), W, Inches(0.4),
         "Happy to ship a working prototype, or install in person.",
         13, W40, italic=True, alignment=PP_ALIGN.CENTER)

# Contact strip
add_text(s, Inches(0), Inches(6.4), W, Inches(0.3),
         "kbarnoski@gmail.com",
         13, PURPLE, font_name="Consolas",
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(6.8), W, Inches(0.25),
         "getresonance.vercel.app",
         11, W30, font_name="Consolas",
         alignment=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "Resonance-Installation-Deck.pptx")
prs.save(out_path)
print(f"Saved → {out_path}")
