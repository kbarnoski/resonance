#!/usr/bin/env python3
"""Generate the Resonance Vision Deck as a PowerPoint file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Constants ──────────────────────────────────────────────────────────
W = Inches(13.333)  # 1440px at 108dpi → widescreen 16:9
H = Inches(7.5)
BG = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)  # canonical app purple — see docs/brand/brand-system.md
MUTED = lambda a: RGBColor(int(255*a), int(255*a), int(255*a))  # grey at opacity

# Approximate opacity-on-black as solid grey
W90 = WHITE
W75 = RGBColor(0xBF, 0xBF, 0xBF)
W60 = RGBColor(0x99, 0x99, 0x99)
W50 = RGBColor(0x80, 0x80, 0x80)
W45 = RGBColor(0x73, 0x73, 0x73)
W40 = RGBColor(0x66, 0x66, 0x66)
W35 = RGBColor(0x59, 0x59, 0x59)
W30 = RGBColor(0x4D, 0x4D, 0x4D)
W25 = RGBColor(0x40, 0x40, 0x40)
W22 = RGBColor(0x38, 0x38, 0x38)
W20 = RGBColor(0x33, 0x33, 0x33)
W15 = RGBColor(0x26, 0x26, 0x26)

P60 = RGBColor(0xC4, 0xB5, 0xFD)  # acc-light — Tailwind violet-300
P50 = RGBColor(0xB3, 0xA0, 0xFB)
P40 = RGBColor(0xA3, 0x8B, 0xF8)
P25 = RGBColor(0x6E, 0x4A, 0xC4)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank layout


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=14,
             color=WHITE, bold=False, italic=False, font_name="Calibri",
             alignment=PP_ALIGN.LEFT, line_spacing=None, spacing_after=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if spacing_after is not None:
        p.space_after = Pt(spacing_after)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_name="Calibri"):
    """lines = list of (text, font_size, color, bold, italic, line_spacing, spacing_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, fs, clr, bld, ital = line[0], line[1], line[2], line[3], line[4]
        ls = line[5] if len(line) > 5 else None
        sa = line[6] if len(line) > 6 else None
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.italic = ital
        p.font.name = font_name
        if ls:
            p.line_spacing = Pt(ls)
        if sa is not None:
            p.space_after = Pt(sa)
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


# ── Margins ────────────────────────────────────────────────────────────
LM = Inches(0.9)   # left margin (~96px)
TM = Inches(0.67)  # top margin (~72px)
CW = Inches(11.5)  # content width


# ══════════════════════════════════════════════════════════════════════
# SLIDE 01 — Title
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Resonance wave mark (simple line as stand-in)
add_text(s, LM, Inches(2.5), Inches(1), Inches(0.4), "⌇", 24, W60, font_name="Calibri")

add_text(s, LM, Inches(3.0), Inches(8), Inches(1.0),
         "Resonance", 52, WHITE, bold=True, font_name="Calibri")

add_text(s, LM, Inches(4.1), Inches(8), Inches(0.8),
         "A personal audio workspace for musicians who want\nto understand, visualize, and experience their music.",
         17, W50, font_name="Calibri", line_spacing=26)

add_text(s, LM, Inches(6.5), Inches(8), Inches(0.3),
         "KAREL BARNOSKI    |    MARCH 2026    |    V1.0",
         9, W25, font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 02 — The Problem
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.4), Inches(3), Inches(0.25),
         "THE PROBLEM", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.75), Inches(10), Inches(1.2),
         "Musicians capture ideas.\nThen forget them.", 48, WHITE, bold=True,
         font_name="Calibri", line_spacing=56)

add_multiline(s, LM, Inches(3.3), Inches(7), Inches(1.8), [
    ("Hundreds of recordings pile up unnamed in phone folders. \"What key was that in?\" "
     "\"What were those chords?\" No way to search, compare, or develop ideas without "
     "sitting back down at the instrument.", 14, W45, False, False, 22, 12),
    ("Creative output goes unanalyzed, undeveloped, lost.", 14, W40, False, False, 22, 16),
    ("The gap between capturing an idea and understanding it is where songs go to die.",
     14, P60, False, False, 22),
])


# ══════════════════════════════════════════════════════════════════════
# SLIDE 03 — The Solution
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(1.0), Inches(3), Inches(0.25),
         "THE SOLUTION", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(1.35), Inches(10), Inches(1.2),
         "Resonance makes every\nrecording self-aware.", 48, WHITE, bold=True,
         font_name="Calibri", line_spacing=56)

add_multiline(s, LM, Inches(2.9), Inches(7), Inches(2.0), [
    ("Upload a musical recording. Resonance transcribes it, detects the key, chords, "
     "tempo, and progressions — then gives you an AI music theory coach to help develop it.",
     14, W45, False, False, 22, 16),
    ("But understanding is only half of the experience. Resonance also transforms your "
     "music into something you can inhabit — immersive visual worlds shaped by the "
     "harmonic DNA of what you've played.", 14, W45, False, False, 22),
])

# Studio / The Room cards
for i, (title, desc) in enumerate([
    ("Studio", "Analyze, study, and understand your music"),
    ("The Room", "An immersive space for music and visuals"),
]):
    x = LM + Inches(i * 3.2)
    add_rect(s, x, Inches(5.1), Inches(2.8), Inches(0.04), fill_color=W15)
    add_text(s, x + Inches(0.1), Inches(5.2), Inches(2.6), Inches(0.3),
             title, 14, WHITE, bold=True)
    add_text(s, x + Inches(0.1), Inches(5.55), Inches(2.6), Inches(0.5),
             desc, 12, W40, line_spacing=18)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 04 — How It Works
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, Inches(2.6), Inches(3), Inches(0.25),
         "HOW IT WORKS", 9, PURPLE, font_name="Consolas")

steps = [
    ("01", "Upload", "Drag and drop your recordings. M4A, MP3, WAV. iPhone voice memos work perfectly."),
    ("02", "Analyze", "AI transcribes notes, detects key, chords, tempo, time signature, and recurring progressions."),
    ("03", "Understand", "Teaching summaries break down sections. Chat with AI about harmony, development, style."),
    ("04", "Experience & Share", "Step into The Room. Your music becomes visual worlds — shaders, poetry, AI imagery. Share a link and anyone can step inside."),
]
card_w = Inches(2.7)
gap = Inches(0.2)
for i, (num, title, desc) in enumerate(steps):
    x = LM + (card_w + gap) * i
    y = Inches(3.2)
    add_rect(s, x, y, card_w, Inches(2.1), border_color=W15)
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.25),
             num, 10, PURPLE, font_name="Consolas")
    add_text(s, x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.4), Inches(0.35),
             title, 18, WHITE, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(1.1), card_w - Inches(0.4), Inches(0.9),
             desc, 11, W40, line_spacing=17)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 05 — Studio
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "STUDIO", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(0.7),
         "Your analytical workspace.", 42, WHITE, bold=True)

features = [
    ("AI Analysis Pipeline", "Client-side note transcription, key detection, chord detection with inversions, tempo, time signature, recurring progressions. MIDI export."),
    ("Teaching Summaries", "AI-generated section breakdowns, harmonic vocabulary, relearning tips. Structured to help you sit back down and play."),
    ("Music Theory Chat", "Per-recording, comparison, and library-wide AI conversations grounded in your analysis data."),
    ("Smart Library", "Search, tag, and organize recordings. Collections, inline editing, full metadata display."),
    ("Insights Dashboard", "Musical DNA, key distribution, chord frequency, cross-recording similarity, AI-generated library portrait."),
    ("Visual Analysis Tools", "Waveform player, chord timeline with synced playhead, piano roll with velocity, time-stamped markers."),
]
cols = 4
for i, (title, desc) in enumerate(features):
    col = i % cols
    row = i // cols
    x = LM + Inches(col * 3.0)
    y = Inches(2.0) + Inches(row * 1.5)
    add_text(s, x, y, Inches(2.7), Inches(0.25), title, 13, WHITE, bold=True)
    add_text(s, x, y + Inches(0.35), Inches(2.7), Inches(0.9), desc, 11, W40, line_spacing=17)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 06 — The Room
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "THE ROOM", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(0.7),
         "Step inside your music.", 42, WHITE, bold=True)

room_features = [
    ("200+ Shader Visualizations", "Six categories — Elemental, Visionary, Cosmic, Organic, Geometry, Dark — plus immersive 3D worlds. Smooth crossfades between modes."),
    ("Journey System", "240+ journeys across 16 realms. Six-phase emotional arcs with multi-layer shader compositing and palette theming."),
    ("AI Poetry & Voice", "Real-time poetic overlays synced to mood and phase. Whisper narration in 12 languages. 35 creative angles that rotate to stay fresh."),
    ("Story Mode", "Six-phase narrative arcs generated from your music's harmonic character. AI-written paragraphs unfold across the journey."),
    ("Custom Journeys", "Create from a story prompt or auto-generate from audio analysis. AI maps your music's mood to a realm and composes the visual arc."),
    ("Shareable Rooms", "Public links to any journey. Anyone can experience your music — shaders, poetry, narrative — without logging in."),
]
for i, (title, desc) in enumerate(room_features):
    col = i % 3
    row = i // 3
    x = LM + Inches(col * 3.9)
    y = Inches(2.0) + Inches(row * 1.4)
    add_text(s, x, y, Inches(3.5), Inches(0.25), title, 13, WHITE, bold=True)
    add_text(s, x, y + Inches(0.35), Inches(3.5), Inches(0.8), desc, 11, W40, line_spacing=17)

# Journey phases at bottom
phases = ["Threshold", "Expansion", "Transcendence", "Illumination", "Return", "Integration"]
add_text(s, LM, Inches(5.85), Inches(3), Inches(0.2),
         "JOURNEY PHASES", 8, W25, font_name="Consolas")
phase_colors = [P40, P50, P60, PURPLE, P50, P40]
for i, (phase, clr) in enumerate(zip(phases, phase_colors)):
    x = LM + Inches(i * 1.95)
    # line
    add_rect(s, x, Inches(6.2), Inches(1.8), Inches(0.02), fill_color=clr)
    add_text(s, x, Inches(6.3), Inches(1.8), Inches(0.2),
             phase, 8, clr, font_name="Consolas", alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 07 — Share the Experience
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "SHARING", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(1.0),
         "Music was never meant\nto be a file.", 48, WHITE, bold=True,
         font_name="Calibri", line_spacing=56)

add_text(s, LM, Inches(2.4), Inches(7), Inches(0.8),
         "When you share a recording today, someone gets an audio file. They press play "
         "in a browser tab. That's it. Resonance turns sharing into an invitation — "
         "step inside this music and experience what I experienced.",
         14, W45, line_spacing=22)

layers = [
    ("LAYER 1", "Your Music", "The recording itself — streamed in full quality through The Room's global audio engine."),
    ("LAYER 2", "Shader Visuals", "Multi-layer WebGL shaders composited in real time — rotating through curated palettes across six journey phases."),
    ("LAYER 3", "AI Poetry", "Poetic text overlays generated from the music's mood. Whispered narration in 12 languages. Fades in and out with the phases."),
    ("LAYER 4", "AI Imagery", "Real-time generated images from the journey's visual vocabulary — landscapes, textures, abstract forms that emerge from your harmonic DNA."),
    ("LAYER 5", "Narrative Arc", "A six-phase story structure — Threshold to Integration — that gives every shared experience a beginning, climax, and resolution."),
]
layer_purples = [P40, P50, P60, PURPLE, PURPLE]
for i, (label, title, desc) in enumerate(layers):
    y = Inches(3.3) + Inches(i * 0.58)
    add_rect(s, LM, y + Inches(0.5), CW, Inches(0.01), fill_color=W15)
    add_text(s, LM, y, Inches(0.8), Inches(0.25),
             label, 8, layer_purples[i], font_name="Consolas")
    add_text(s, LM + Inches(1.2), y, Inches(1.5), Inches(0.25),
             title, 13, WHITE, bold=True)
    add_text(s, LM + Inches(3.0), y, Inches(8.5), Inches(0.5),
             desc, 11, W40, line_spacing=17)

add_text(s, LM, Inches(6.7), Inches(8), Inches(0.3),
         "One link. Five layers. A complete experience anyone can open.",
         12, W50, bold=True)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 08 — Installation & Performance
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "INSTALLATION MODE", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(0.7),
         "Not just software. A medium.", 42, WHITE, bold=True)

add_text(s, LM, Inches(1.75), Inches(7), Inches(0.6),
         "A dedicated kiosk mode designed for physical spaces — galleries, lobbies, "
         "performance venues, studios, and public environments.",
         14, W45, line_spacing=22)

use_cases = [
    ("Gallery Installation", "A pianist's recordings visualized as living art. Visitors walk into a room where music creates its own visual world."),
    ("Live Performance", "A musician performs while The Room generates real-time visuals behind them. A unique visual experience for every show."),
    ("Meditation & Healing", "Spa lobbies, therapy offices, yoga studios. Ambient music with gently evolving visuals and whispered poetry."),
    ("Retail & Hospitality", "Hotels, restaurants, boutiques. Curated sonic-visual environments that set a mood without demanding attention."),
    ("Museum & Cultural", "Pair historical recordings with visual journeys. A Chopin nocturne through the Winter realm. Jazz through Cosmos."),
    ("Studio Waiting Rooms", "Recording studios, music schools, creative agencies. Showcase original music as a visual experience while clients wait."),
]
for i, (title, desc) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = LM + Inches(col * 3.9)
    y = Inches(2.7) + Inches(row * 1.6)
    add_rect(s, x, y, Inches(3.5), Inches(1.3), border_color=RGBColor(0x2A, 0x1A, 0x4A))
    add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(3.1), Inches(0.25),
             title, 12, WHITE, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(3.1), Inches(0.7),
             desc, 11, W40, line_spacing=17)

# Tech specs footer
specs = "Auto-play curated queue   |   Fullscreen kiosk mode   |   Cursor auto-hide   |   Loop indefinitely   |   URL-configurable"
add_text(s, LM, Inches(6.7), CW, Inches(0.25),
         specs, 8, W25, font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════
# SLIDE 09 — Insights
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "INSIGHTS", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(0.7),
         "Your musical DNA.", 42, WHITE, bold=True)

add_text(s, LM, Inches(1.75), Inches(5), Inches(0.4),
         "See patterns across your entire creative output.", 14, W45)

insights = [
    ("Key Distribution", "Which keys you gravitate toward. Major vs minor tendencies. Patterns you didn't know you had."),
    ("Chord Vocabulary", "Frequency analysis of every chord across your library. Jazz-influenced, diatonic, chromatic — see your harmonic fingerprint."),
    ("Recurring Progressions", "Identify chord sequences that appear across multiple recordings. Your signature moves surfaced automatically."),
    ("Recording Similarity", "Pairwise comparison by key, tempo, and chord overlap. Discover which of your recordings are harmonic siblings."),
    ("AI Library Portrait", "Claude analyzes your entire library and writes a portrait of your musical identity — clusters, standout pieces, development arcs."),
    ("Compare Mode", "Side-by-side analysis of any two recordings with match/differs badges and dedicated AI comparison chat."),
]
for i, (title, desc) in enumerate(insights):
    col = i % 2
    row = i // 2
    x = LM + Inches(col * 6.0)
    y = Inches(2.6) + Inches(row * 1.3)
    add_text(s, x, y, Inches(5.5), Inches(0.25), title, 13, WHITE, bold=True)
    add_text(s, x, y + Inches(0.35), Inches(5.5), Inches(0.7), desc, 11, W40, line_spacing=17)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Technology
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "TECHNOLOGY", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(1.0),
         "Built for musicians,\npowered by AI.", 42, WHITE, bold=True,
         line_spacing=50)

tech_left = [
    ("Frontend", "Next.js 15, React 19, TypeScript, Tailwind v4"),
    ("AI", "Claude (Anthropic) via Vercel AI SDK"),
    ("Audio Intelligence", "Spotify Basic Pitch, tonal.js"),
    ("Visualization", "WebGL (200+ shaders), Three.js (3D)"),
]
tech_right = [
    ("Audio Engine", "Web Audio API, WaveSurfer.js"),
    ("Image Generation", "FAL Flux (real-time AI imagery)"),
    ("Backend", "Supabase (Auth, PostgreSQL, Storage)"),
    ("Hosting", "Vercel (serverless, edge)"),
]

for col, items in enumerate([tech_left, tech_right]):
    x_label = LM + Inches(col * 6.0)
    x_val = x_label + Inches(1.8)
    for i, (label, val) in enumerate(items):
        y = Inches(2.6) + Inches(i * 0.55)
        add_text(s, x_label, y, Inches(1.7), Inches(0.25), label, 11, W35)
        add_text(s, x_val, y, Inches(4.0), Inches(0.25), val, 10, W60, font_name="Consolas")
        add_rect(s, x_label, y + Inches(0.35), Inches(5.5), Inches(0.01), fill_color=W15)

# Architecture highlights
arch = [
    ("Global audio singleton", "Survives route changes — seamless navigation between Studio and The Room"),
    ("Client-side analysis", "No audio sent to servers — TensorFlow.js runs note transcription in the browser"),
    ("Independent journey engine", "Phase transitions, shader rotation, and timing run outside React for smooth performance"),
]
for i, (title, desc) in enumerate(arch):
    x = LM + Inches(i * 3.9)
    y = Inches(5.6)
    add_rect(s, x, y, Inches(3.5), Inches(1.0), fill_color=RGBColor(0x12, 0x0E, 0x20))
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.1), Inches(0.25),
             title, 10, W60, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.45), Inches(3.1), Inches(0.5),
             desc, 9, W30, line_spacing=14)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Roadmap
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "ROADMAP", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(0.7),
         "Where we're going.", 42, WHITE, bold=True)

roadmap = [
    ("Q1 2026", "SHIPPED", PURPLE, WHITE,
     "Foundation + The Room", [
         "Full analysis pipeline (key, chords, tempo, MIDI)",
         "AI teaching summaries and chat",
         "200+ shaders, 240+ journeys, 16 realms",
         "Poetry, voice, story overlays",
         "Installation mode, public sharing",
     ]),
    ("Q2 2026", "NEXT", P60, W75,
     "Intelligence & Generation", [
         "AI journey auto-generation from analysis",
         "Real-time AI imagery during journeys",
         "In-browser recording",
         "Harmonic search",
         "Practice mode with looping + metronome",
     ]),
    ("Q3–Q4 2026", "LATER", P40, W50,
     "Platform & Social", [
         "Public profiles and journey discovery",
         "Collaborative listening sessions",
         "Multi-instrument support",
         "Version tracking for ideas",
         "Mobile app (PWA)",
     ]),
    ("2027+", "FUTURE", P25, W40,
     "Ecosystem", [
         "Live transcription as you play",
         "AI arrangement and song builder",
         "Style DNA matching",
         "Education platform",
         "DAW export, developer API",
     ]),
]
for i, (quarter, status, q_color, title_color, title, items) in enumerate(roadmap):
    col = i % 2
    row = i // 2
    x = LM + Inches(col * 6.0)
    y = Inches(2.0) + Inches(row * 2.5)

    add_text(s, x, y, Inches(1.5), Inches(0.2),
             quarter, 9, q_color, font_name="Consolas")
    add_text(s, x + Inches(1.6), y, Inches(1), Inches(0.2),
             status, 9, W25)
    add_text(s, x, y + Inches(0.35), Inches(5), Inches(0.3),
             title, 14, title_color, bold=True)

    # Dimming based on timeline position
    item_colors = [W35, W30, W25, W22]
    for j, item in enumerate(items):
        add_text(s, x, y + Inches(0.75) + Inches(j * 0.28), Inches(5), Inches(0.25),
                 item, 10, item_colors[min(i, 3)], line_spacing=15)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Vision
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "NORTH STAR", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(1.2),
         "The musical brain\nyou've always wanted.", 52, WHITE, bold=True,
         font_name="Calibri", line_spacing=62)

add_multiline(s, LM, Inches(2.7), Inches(7.5), Inches(2.5), [
    ("Imagine picking up your instrument. You play something — 30 seconds, two minutes, "
     "whatever comes out. Before you even set it down, the chords appear on screen in "
     "real time. The app recognizes this is a variation of something you played three months ago.",
     13, W45, False, False, 21, 16),
    ("And when you're done — you step into The Room. Your music becomes a world. "
     "Shaders paint the harmonic structure in light. Poetry whispers what the chords are "
     "saying. AI-generated imagery emerges from the mood of what you just played. "
     "You share the link. A friend opens it on their phone. They're inside your music.",
     13, W45, False, False, 21),
])

add_text(s, LM, Inches(5.4), Inches(10), Inches(0.8),
         "Resonance transforms how musicians share their work. Not as a file. Not as a stream. "
         "As an experience — visual, poetic, immersive — that anyone can step into with a single link.",
         15, W60, bold=True, line_spacing=24)

add_text(s, LM, Inches(6.4), Inches(10), Inches(0.6),
         "Not a DAW. Not a notation app. Not a screensaver. A creative partner that turns "
         "every idea into something you can understand, inhabit, and share with the world.",
         12, W30, italic=True, line_spacing=20)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 13 — About
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, LM, TM, Inches(3), Inches(0.25),
         "ABOUT", 9, PURPLE, font_name="Consolas")

add_text(s, LM, Inches(0.95), Inches(10), Inches(0.7),
         "Karel Barnoski", 42, WHITE, bold=True)

bio_left = [
    ("Design Director at Workday. 15+ years leading UX for enterprise products at "
     "Workday, GE, and Kodak. BFA Illustration and MFA Computer Graphics from Rochester "
     "Institute of Technology, with classical training from Barnstone Studios.",
     12, W50, False, False, 19, 14),
    ("Founder of 2octave, an audio design studio specializing in product sonification — "
     "creating award-winning sounds for products ranging from digital cameras to kitchen "
     "appliances. Published in UX Magazine on the science of sound in product design.",
     12, W50, False, False, 19),
]
add_multiline(s, LM, Inches(2.0), Inches(5.5), Inches(3.0), bio_left)

bio_right = [
    ("Pianist and composer. Released \"Welcome Home,\" a solo piano album self-produced "
     "during quarantine. Music on Spotify and SoundCloud.",
     12, W50, False, False, 19, 14),
    ("Resonance exists because Karel is both the designer and the user — a musician who "
     "records hundreds of ideas and needed a better way to understand and develop them. "
     "The Room exists because understanding music was never enough. He wanted to live inside it.",
     12, W50, False, False, 19),
]
add_multiline(s, LM + Inches(6.2), Inches(2.0), Inches(5.5), Inches(3.0), bio_right)

add_text(s, LM, Inches(6.3), Inches(7), Inches(0.6),
         "\"Sound is perceived very much in the same way as any other stimuli. "
         "By manipulating sound, you can affect the user's perception of an experience.\"",
         12, W25, italic=True, line_spacing=19)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Experience It (CTA)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

add_text(s, Inches(0), Inches(2.6), W, Inches(1.0),
         "Resonance", 60, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(3.7), W, Inches(0.5),
         "Analyze. Visualize. Experience. Share.", 16, W40,
         alignment=PP_ALIGN.CENTER, line_spacing=24)

add_text(s, Inches(0), Inches(5.8), W, Inches(0.3),
         "getresonance.vercel.app", 12, PURPLE, font_name="Consolas",
         alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(6.2), W, Inches(0.3),
         "Karel Barnoski  ·  March 2026", 10, W25,
         alignment=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "Resonance-Vision-Deck.pptx")
prs.save(out_path)
print(f"Saved → {out_path}")
